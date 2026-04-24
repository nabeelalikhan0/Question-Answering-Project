import json
import os
from uuid import uuid4
import PyPDF2
import docx
from django.shortcuts import render,redirect,HttpResponse
from django.core.mail import send_mail
from . import models, forms
from .gemini import ask_gemini
from django.contrib.auth.forms import UserCreationForm
from .forms import StyledSignupForm, StyledLoginForm
from django.contrib.auth.views import LoginView
from django.contrib import messages
from django.db.models import Max,Q
from django.urls import reverse
from django.contrib.auth.decorators import login_required
from django.shortcuts import render, redirect
from django.contrib import messages
from django.contrib.auth import update_session_auth_hash
from django.contrib.auth.forms import PasswordChangeForm
from django.contrib.auth import get_user_model



# Home Page
def send_email(email,subject,message):
    recipient_list = [email]
    send_mail(
        subject=subject,
        message=message,
        from_email='nabeelalikhan314@gmail.com',
        recipient_list=recipient_list,
        fail_silently=False,
    )
    return HttpResponse("Message Sent! Successfully")


def index(request):
    return render(request, "index.html")

def about(request):
    return render(request,"about.html")

# Contact Form
def contact(request):
    if request.method == "POST":
        form = forms.ContactForm(request.POST)
        if form.is_valid():
            form.save()
            messages.success(request, "✅ Your message has been sent successfully!")
            email = form.cleaned_data.get("email")  # get email value
            
            try:
                send_email(
                    email,
                    "Thank You for Contacting RAGQA Support",
                    "Hello,\n\nThank you for reaching out to RAGQA Support. "
                    "We have received your message and our team will respond to you as soon as possible.\n\n"
                    "Best regards,\nThe RAGQA Support Team"
                )

            except Exception as e:
                messages.error(request, f"⚠️ Email sending failed: {e}")

                
            return redirect('contact')  # Redirect to clear POST data


    else:
        form = forms.ContactForm()
    return render(request, "contact.html", {"form": form})


User = get_user_model()

@login_required
def profile_view(request):
    return render(request, "profile.html")

@login_required
def edit_profile(request):
    if request.method == "POST":
        first_name = request.POST.get("first_name")
        last_name = request.POST.get("last_name")
        email = request.POST.get("email")

        user = request.user
        user.first_name = first_name
        user.last_name = last_name
        user.email = email
        user.save()

        messages.success(request, "✅ Your profile has been updated successfully.")
        return redirect("profile")
    return redirect("profile")

@login_required
def change_password(request):
    if request.method == "POST":
        form = PasswordChangeForm(request.user, request.POST)
        if form.is_valid():
            user = form.save()
            update_session_auth_hash(request, user)  # Keep user logged in
            messages.success(request, "✅ Your password has been updated successfully.")
            return redirect("profile")
        else:
            for field, errors in form.errors.items():
                for error in errors:
                    messages.error(request, f"{field}: {error}")
    else:
        form = PasswordChangeForm(request.user)

    return render(request, "profile.html", {"password_form": form})



# Chatbot (file upload + chat in one view)

def _extract_text_from_any_file(file_path: str, ext: str, content_type: str | None = None) -> str:
    """
    Best-effort text extraction for many file types.
    Falls back to decoding bytes if no dedicated extractor fits.
    """
    ext = (ext or "").lower()

    # 1) Plain text
    if ext in (".txt", ".log", ".md", ".csv"):
        # CSV shown as plain text; you can switch to pandas if you prefer a table text
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    # 2) PDF
    if ext == ".pdf":
        text_chunks = []
        with open(file_path, "rb") as pdff:
            reader = PyPDF2.PdfReader(pdff)
            for page in reader.pages:
                t = page.extract_text() or ""
                if t:
                    text_chunks.append(t)
        return "\n".join(text_chunks).strip() or "[No text extracted from PDF]"

    # 3) DOCX
    if ext == ".docx":
        d = docx.Document(file_path)
        return "\n".join(p.text for p in d.paragraphs).strip() or "[No text extracted from DOCX]"

    # 4) PPTX
    if ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            out = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        out.append(shape.text)
            return "\n".join(out).strip() or "[No text in PPTX]"
        except Exception as e:
            return f"[PPTX read error: {e}]"

    # 5) XLSX/XLS (first ~1000 rows)
    if ext in (".xlsx", ".xls"):
        try:
            import pandas as pd
            df = pd.read_excel(file_path, nrows=1000)  # limit for sanity
            return df.to_string(index=False)
        except Exception as e:
            return f"[Excel read error: {e}]"

    # 6) HTML
    if ext in (".html", ".htm"):
        try:
            from bs4 import BeautifulSoup
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                soup = BeautifulSoup(f, "html.parser")
            return soup.get_text(separator=" ", strip=True)
        except Exception as e:
            return f"[HTML parse error: {e}]"

    # 7) JSON
    if ext == ".json":
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                data = json.load(f)
            return json.dumps(data, indent=2, ensure_ascii=False)
        except Exception:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()

    # 8) Images -> OCR (best effort)
    if ext in (".png", ".jpg", ".jpeg", ".bmp", ".webp", ".tif", ".tiff"):
        try:
            import pytesseract
            from PIL import Image
            text = pytesseract.image_to_string(Image.open(file_path))
            text = (text or "").strip()
            return text if text else "[No text detected via OCR]"
        except Exception as e:
            return f"[OCR unavailable: {e}]"

    # 9) Fallback: try decode bytes
    try:
        with open(file_path, "rb") as f:
            raw = f.read()
        try:
            return raw.decode("utf-8")
        except UnicodeDecodeError:
            try:
                return raw.decode("latin-1")
            except UnicodeDecodeError:
                return "[Unsupported file type for text extraction]"
    except Exception as e:
        return f"[File read error: {e}]"

# ---------- View ----------


def chatbot(request):
    # Ensure session exists
    if not request.session.session_key:
        request.session.create()

    # Start new chat
    if request.GET.get("new_chat") == "1":
        new_id = uuid4().hex
        request.session["active_chat_session"] = new_id
        return redirect(f"{reverse('chatbot')}?session_id={new_id}")

    # Active session
    selected_session_id = (
        request.GET.get("session_id")
        or request.session.get("active_chat_session")
        or request.session.session_key
    )
    request.session["active_chat_session"] = selected_session_id

    file_text = None
    ai_response = None

    # POST: File upload or chat message
    if request.method == "POST":
        # File upload
        if "file" in request.FILES:
            up = request.FILES["file"]
            obj = models.PreprocessText.objects.create(
                session_id=selected_session_id,
                file=up,
            )
            ext = os.path.splitext(obj.file.name)[1].lower()
            content_type = getattr(up, "content_type", None)
            extracted = _extract_text_from_any_file(obj.file.path, ext, content_type)
            obj.file_text = extracted
            obj.save()

            if not extracted or (extracted.startswith("[") and "error" in extracted.lower()):
                messages.warning(
                    request,
                    "Uploaded, but text extraction may be limited. You can still ask questions."
                )
            return redirect(f"{reverse('chatbot')}?session_id={selected_session_id}")

        # Chat message
        user_message = (request.POST.get("message") or "").strip()
        if user_message:
            last_file = models.PreprocessText.objects.filter(session_id=selected_session_id).order_by("-id").first()
            file_text = last_file.file_text if last_file else None
            if not file_text:
                messages.error(request, "Please upload a file first.")
            else:
                from .gemini import ask_gemini
                ai_response = ask_gemini(f"{file_text}\n\n{user_message}")
                models.ChatHistory.objects.create(
                    user=request.user if request.user.is_authenticated else None,
                    session_id=selected_session_id,
                    user_message=user_message,
                    ai_response=ai_response,
                )
            return redirect(f"{reverse('chatbot')}?session_id={selected_session_id}")

    # GET: page load
    last_file = models.PreprocessText.objects.filter(session_id=selected_session_id).order_by("-id").first()
    if last_file:
        file_text = last_file.file_text

    # Chat search
    search_query = (request.GET.get("search") or "").strip()
    if search_query and request.user.is_authenticated:
        # Search across all sessions for this user
        chat_qs = models.ChatHistory.objects.filter(user=request.user)
        chat_qs = chat_qs.filter(
            Q(user_message__icontains=search_query) | Q(ai_response__icontains=search_query)
        )
    else:
        # Default: show current session
        chat_qs = models.ChatHistory.objects.filter(session_id=selected_session_id)
    chat_history = chat_qs.order_by("created_at")

    # Previous sessions for user
    previous_sessions = []
    if request.user.is_authenticated:
        previous_sessions = (
            models.ChatHistory.objects
            .filter(user=request.user)
            .values("session_id")
            .annotate(last_message_time=Max("created_at"))
            .order_by("-last_message_time")
        )

    return render(request, "chatbot.html", {
        "file_text": file_text,
        "ai_response": ai_response,
        "chat_history": chat_history,
        "previous_sessions": previous_sessions,
        "selected_session_id": selected_session_id,
        "search_query": search_query,
    })


def signup(request):
    if request.method == "POST":
        form = StyledSignupForm(request.POST)
        if form.is_valid():
            form.save()
            return redirect('login')
    else:
        form = StyledSignupForm()
    return render(request, 'registration/signup.html', {'form': form})


class CustomLoginView(LoginView):
    authentication_form = StyledLoginForm
    template_name = "registration/login.html"


# views.py

def subscribe(request):
    if request.method == "POST":
        email = request.POST.get("email") or (request.user.email if request.user.is_authenticated else None)

        if not email:
            messages.error(request, "Please enter a valid email.")
            return redirect(request.META.get("HTTP_REFERER", "/"))

        # 🟢 Improved logic for creating a subscriber
        user_instance = request.user if request.user.is_authenticated else None
        
        # Use email as the main lookup and provide user as a default
        obj, created = models.subscribers.objects.update_or_create(
            email=email,
            defaults={'user': user_instance}
        )

        # Send confirmation email (your existing logic is fine here)
        subject = "RAGQA Subscription Confirmation"
        message = (
            f"Dear {user_instance.first_name if user_instance and user_instance.first_name else 'Subscriber'},\n\n"
            "Thank you for subscribing to RAGQA.\n"
            "You’re now part of our intelligent Q&A community where knowledge meets precision.\n\n"
            "We look forward to helping you find accurate answers instantly.\n\n"
            "Best regards,\n"
            "The RAGQA Team"
        )
        send_email(email, subject, message)

        if created:
            messages.success(request, "Subscription successful! A confirmation email has been sent.")
        else:
            messages.info(request, "You are already subscribed. We've resent your confirmation email.")

        return render(request, "subscribe.html", {"email": email})
    
    return render(request, "subscribe_failed.html")


from django.views.decorators.csrf import csrf_exempt
from django.http import JsonResponse

@csrf_exempt
def api_chat(request):
    if request.method == "POST":
        selected_session_id = request.POST.get("session_id") or "default_api_session"
        user_message = request.POST.get("message", "").strip()
        
        file_text = None
        
        # Handle file upload if present
        if "file" in request.FILES:
            up = request.FILES["file"]
            obj = models.PreprocessText.objects.create(
                session_id=selected_session_id,
                file=up,
            )
            ext = os.path.splitext(obj.file.name)[1].lower()
            extracted = _extract_text_from_any_file(obj.file.path, ext)
            obj.file_text = extracted
            obj.save()
            file_text = extracted
        else:
            # Fallback to last uploaded file for this session
            last_file = models.PreprocessText.objects.filter(session_id=selected_session_id).order_by("-id").first()
            if last_file:
                file_text = last_file.file_text

        if user_message:
            if not file_text:
                return JsonResponse({"error": "Please upload a file first."}, status=400)
            
            ai_response = ask_gemini(f"{file_text}\n\n{user_message}")
            
            user_id = request.POST.get("user_id")
            user_obj = request.user if request.user.is_authenticated else User.objects.filter(id=user_id).first()

            models.ChatHistory.objects.create(
                user=user_obj,
                session_id=selected_session_id,
                user_message=user_message,
                ai_response=ai_response,
            )
            
            return JsonResponse({
                "role": "ai",
                "content": ai_response,
                "session_id": selected_session_id
            })
            
        return JsonResponse({"error": "No message provided"}, status=400)

from django.contrib.auth import authenticate, login, logout
from django.core.validators import validate_email
from django.core.exceptions import ValidationError

@csrf_exempt
def api_signup(request):
    if request.method == "POST":
        try:
            data = json.loads(request.body) if request.content_type == "application/json" else request.POST
            username = data.get("username")
            email = data.get("email")
            password = data.get("password")
            first_name = data.get("first_name", "")
            
            if not username or not password or not email:
                return JsonResponse({"error": "Missing fields"}, status=400)
            
            if User.objects.filter(username=username).exists():
                return JsonResponse({"error": "Username already exists"}, status=400)
            
            user = User.objects.create_user(username=username, email=email, password=password, first_name=first_name)
            return JsonResponse({"message": "Account created successfully", "user_id": user.id})
        except Exception as e:
            return JsonResponse({"error": str(e)}, status=500)
    return JsonResponse({"error": "POST required"}, status=405)

@csrf_exempt
def api_login(request):
    if request.method == "POST":
        try:
            data = json.loads(request.body) if request.content_type == "application/json" else request.POST
            username = data.get("username")
            password = data.get("password")
            
            user = authenticate(request, username=username, password=password)
            if user:
                login(request, user)
                return JsonResponse({
                    "message": "Login successful",
                    "user": {
                        "id": user.id,
                        "username": user.username,
                        "email": user.email,
                        "first_name": user.first_name
                    }
                })
            return JsonResponse({"error": "Invalid credentials"}, status=401)
        except Exception as e:
            return JsonResponse({"error": str(e)}, status=500)
    return JsonResponse({"error": "POST required"}, status=405)

@csrf_exempt
def api_logout(request):
    logout(request)
    return JsonResponse({"message": "Logged out successfully"})

@csrf_exempt
def api_profile(request):
    user_id = request.GET.get("user_id")
    if not request.user.is_authenticated and not user_id:
        return JsonResponse({"error": "Unauthorized"}, status=401)
    
    user = request.user if request.user.is_authenticated else User.objects.filter(id=user_id).first()
    if not user:
        return JsonResponse({"error": "User not found"}, status=404)
    
    if request.method == "GET":
        return JsonResponse({
            "id": user.id,
            "username": user.username,
            "email": user.email,
            "first_name": user.first_name,
            "last_name": user.last_name,
        })
    
    if request.method == "POST":
        data = json.loads(request.body) if request.content_type == "application/json" else request.POST
        user_id = data.get("user_id")
        user_obj = request.user if request.user.is_authenticated else User.objects.filter(id=user_id).first()
        
        if not user_obj:
            return JsonResponse({"error": "User not found"}, status=404)

        user_obj.first_name = data.get("first_name", user_obj.first_name)
        user_obj.last_name = data.get("last_name", user_obj.last_name)
        user_obj.email = data.get("email", user_obj.email)
        user_obj.save()
        return JsonResponse({"message": "Profile updated successfully"})

@csrf_exempt
def api_contact(request):
    if request.method == "POST":
        data = json.loads(request.body) if request.content_type == "application/json" else request.POST
        name = data.get("name")
        email = data.get("email")
        subject = data.get("subject")
        message = data.get("message")
        
        contact_obj = models.Contact.objects.create(name=name, email=email, subject=subject, message=message)
        
        # Send email (reuse existing logic)
        try:
            send_email(email, f"Contact: {subject}", f"Thank you for contacting us, {name}.\n\nMessage: {message}")
        except:
            pass
            
        return JsonResponse({"message": "Message sent successfully"})
    return JsonResponse({"error": "POST required"}, status=405)

@csrf_exempt
def api_subscribe(request):
    if request.method == "POST":
        try:
            data = json.loads(request.body) if request.content_type == "application/json" else request.POST
            email = data.get("email")
            if not email:
                return JsonResponse({"error": "Email required"}, status=400)
            
            user_instance = request.user if request.user.is_authenticated else None
            obj, created = models.subscribers.objects.update_or_create(
                email=email,
                defaults={'user': user_instance}
            )
            
            # Send confirmation
            send_email(email, "Subscription Confirmed", "Thank you for subscribing to RAG Q&A.")
            
            return JsonResponse({"message": "Subscribed successfully", "created": created})
        except Exception as e:
            return JsonResponse({"error": str(e)}, status=500)
    return JsonResponse({"error": "POST required"}, status=405)

@csrf_exempt
def api_sessions(request):
    user_id = request.GET.get("user_id")
    if not request.user.is_authenticated and not user_id:
        return JsonResponse({"error": "Unauthorized"}, status=401)
    
    user = request.user if request.user.is_authenticated else User.objects.filter(id=user_id).first()
    if not user:
        return JsonResponse({"error": "User not found"}, status=404)

    histories = models.ChatHistory.objects.filter(user=user).order_by('-created_at')
    sessions = []
    seen_sessions = set()
    for h in histories:
        sid = h.session_id or "default"
        if sid not in seen_sessions:
            sessions.append({
                "id": sid,
                "title": h.user_message[:30] + ("..." if len(h.user_message) > 30 else ""),
                "lastMessage": h.ai_response[:50] + ("..." if len(h.ai_response) > 50 else ""),
                "date": h.created_at.isoformat()
            })
            seen_sessions.add(sid)
    
    return JsonResponse({"sessions": sessions})


